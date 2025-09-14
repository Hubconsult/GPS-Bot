import io
from datetime import datetime
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

from font_data import ensure_font


def make_pdf(text: str) -> bytes:
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=A4)

    font_path = ensure_font()
    pdfmetrics.registerFont(TTFont("DejaVu", font_path))
    c.setFont("DejaVu", 12)
    c.setTitle("Document")
    width, height = A4
    y = height - 50
    for line in text.splitlines() or ["(empty)"]:
        c.drawString(40, y, line[:120])  # простая обрезка
        y -= 18
        if y < 40:
            c.showPage()
            c.setFont("DejaVu", 12)
            y = height - 50
    c.showPage()
    c.save()
    return buf.getvalue()

def make_excel(csv_like_text: str) -> bytes:
    """
    Ожидает текст наподобие CSV:
    Заголовок1, Заголовок2
    Значение1, Значение2
    ...
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row in csv_like_text.splitlines():
        cells = [cell.strip() for cell in row.split(",")]
        ws.append(cells)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()

def make_pptx(title_and_bullets: str) -> bytes:
    """
    Текст формата:
    Заголовок: Моя презентация
    - Пункт 1
    - Пункт 2
    ===
    Заголовок: Второй слайд
    - Пункт A
    - Пункт B
    """
    prs = Presentation()
    blocks = [b.strip() for b in title_and_bullets.split("===") if b.strip()]
    if not blocks:
        blocks = [f"Заголовок: Авто-слайды\n- Слайд создан {datetime.now():%Y-%m-%d %H:%M}"]

    for block in blocks:
        lines = [ln.strip() for ln in block.splitlines() if ln.strip()]
        title_line = lines[0] if lines else "Слайд"
        title_text = title_line.replace("Заголовок:", "").strip() if "Заголовок:" in title_line else title_line

        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text

        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = [ln[1:].strip() for ln in lines[1:] if ln.startswith("-")]
        if bullets:
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
        else:
            body.text = " "

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
